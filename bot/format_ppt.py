import os

import pptx
import pptx.util
from pptx.dml.color import RGBColor


class Format:
    def __init__(self, fontBold, fontColor, fontName, fontSize, shapeWidth, shapeHeight):
        self.fontBold = fontBold
        self.fontColor = pptx.dml.color.RGBColor(fontColor[0], fontColor[1], fontColor[2])
        self.fontName = fontName
        self.fontSize = pptx.util.Pt(fontSize)
        self.shapeWidth = pptx.util.Inches(shapeWidth)
        self.shapeHeight = pptx.util.Inches(shapeHeight)


slideHeight = 7.5
slideWidth = 13.333
headingHeightFraction = 1 / 5.5

bodyFormat = Format(
    fontBold=False,
    fontColor=(0, 0, 0),
    fontName='Calibri',
    fontSize=12,
    shapeWidth=slideWidth / 2,
    shapeHeight=slideHeight * (1 - headingHeightFraction),
)
bodyHeadingFormat = Format(
    fontBold=False,
    fontColor=(0, 0, 0),
    fontName='Calibri Light',
    fontSize=16,
    shapeWidth=slideWidth / 2,
    shapeHeight=slideHeight * (1 - headingHeightFraction),
)
captionFormat = Format(
    fontBold=False,
    fontColor=(0, 0, 0),
    fontName='Calibri',
    fontSize=12,
    shapeWidth=slideWidth / 2,
    shapeHeight=slideHeight / 2,
)
headingFormat = Format(
    fontBold=False,
    fontColor=(255, 255, 255),
    fontName='Calibri Light',
    fontSize=14,
    shapeWidth=slideWidth / 2,
    shapeHeight=slideHeight * (headingHeightFraction),
)
subheadingFormat = Format(
    fontBold=False,
    fontColor=(255, 255, 255),
    fontName='Calibri Light',
    fontSize=16,
    shapeWidth=slideWidth / 2,
    shapeHeight=slideHeight * (headingHeightFraction),
)


def detect_title(shape):
    height = (1 / 4) * 9
    width = (1 / 2) * 16
    ratio = height / width
    shapeRatio = shape.height / shape.width
    if (shapeRatio < ratio) and (shape.top <= slideHeight * (1.5 * headingHeightFraction)):
        return detect_text(shape)
    return False


def detect_caption(shape):
    if (shape.height < captionFormat.shapeHeight) and (shape.width < captionFormat.shapeWidth):
        return detect_text(shape)
    return False


def detect_body(shape):
    return detect_text(shape)


def detect_text(shape):
    try:
        if (len(shape.text_frame.text.strip()) > 0):
            return True
        return False
    except:
        return False


def format_body(shape):
    for paragraph in shape.text_frame.paragraphs:
        update_text_format(paragraph, bodyFormat)
        for run in paragraph.runs:
            update_text_format(run, bodyFormat)
    format_shape(shape, bodyFormat)


def format_caption(shape):
    for paragraph in shape.text_frame.paragraphs:
        update_text_format(paragraph, captionFormat)
        for run in paragraph.runs:
            update_text_format(run, captionFormat)


def format_shape(shape, format):
    shape.width = format.shapeWidth
    shape.height = format.shapeHeight


def format_shapes(shapes):
    titleIndex = 0
    bodyIndex = 0
    for shape in shapes:
        try:
            if detect_caption(shape):
                format_caption(shape)
            elif detect_title(shape):
                format_title(shape)
                position_title(shape, titleIndex)
                titleIndex += 1
            elif detect_body(shape):
                format_body(shape)
                position_body(shape, bodyIndex)
                bodyIndex += 1
        except Exception as e:
            print('Shape: {}'.format(e))


def format_title(shape):
    for paragraph in shape.text_frame.paragraphs:
        runs = paragraph.runs
        for run, i in zip(runs, range(len(runs))):
            if i < 1:
                update_text_format(run, headingFormat)
            else:
                update_text_format(run, subheadingFormat)
    format_shape(shape, headingFormat)


def position_body(shape, bodyIndex):
    shape.left = bodyIndex * bodyFormat.shapeWidth
    shape.top = headingFormat.shapeHeight


def position_title(shape, titleIndex):
    shape.left = round(titleIndex * slideWidth / 2)
    shape.top = 0


def format_powerpoint(path):
    ppt = pptx.Presentation(path)
    slides = ppt.slides
    try:
        for slide in slides:
            format_shapes(slide.shapes)
    except Exception as e:
        print('Slide: {}'.format(e))
    ppt.save(path)


def update_text_format(text, format):
    text.font.underline = False
    text.font.bold = format.fontBold
    text.font.name = format.fontName
    text.font.size = format.fontSize
    text.font.color.rgb = format.fontColor
    text.alignment = pptx.enum.text.PP_ALIGN.LEFT


def brute_format_all_powerpoints(path):
    for root, directories, files in os.walk(path):
        try:
            for filename in files:
                filepath = os.path.join(root, filename)
                format_powerpoint(filepath)
        except:
            pass
